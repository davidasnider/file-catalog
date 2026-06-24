import logging
import asyncio
import shutil
from typing import Dict, Any
import pdfplumber
import pytesseract
from PIL import Image
import xlrd
from hachoir.parser import createParser
from hachoir.metadata import extractMetadata

from src.core.plugin_registry import AnalyzerBase, register_analyzer
from src.core.config import config
from src.core.analyzer_names import TEXT_EXTRACTOR_NAME

logger = logging.getLogger(__name__)

SUPPORTED_IMAGE_TYPES = {
    "image/jpeg",
    "image/png",
    "image/bmp",
    "image/tiff",
    "image/gif",
    "image/webp",
    "image/vnd.adobe.photoshop",
}


@register_analyzer(name=TEXT_EXTRACTOR_NAME, depends_on=[], version="1.10")
class TextExtractorPlugin(AnalyzerBase):
    """
    Extracts raw text from common document types (PDFs, docs) and images (OCR).
    """

    def should_run(
        self, file_path: str, mime_type: str, context: Dict[str, Any]
    ) -> bool:
        # Audio and Video are handled by specialized transcribers/analyzers
        if mime_type.startswith(("audio/", "video/")):
            return False

        if config.use_document_ai:
            supported_docai_prefixes = {
                "application/pdf",
                "image/",
                "application/vnd.openxmlformats-officedocument",
            }
            if any(mime_type.startswith(p) for p in supported_docai_prefixes):
                logger.debug(
                    f"Skipping {TEXT_EXTRACTOR_NAME} for {file_path} because Document AI is enabled."
                )
                return False
        return True

    async def analyze(
        self, file_path: str, mime_type: str, context: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Extract text based on the detected MIME type."""
        logger.info(f"Extracting text for {file_path}")

        extracted_text = ""

        try:
            if mime_type == "text/plain":

                def _read_text():
                    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                        return f.read()

                extracted_text = await asyncio.to_thread(_read_text)
            elif mime_type == "text/html":
                from bs4 import BeautifulSoup

                def _read_html():
                    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                        soup = BeautifulSoup(f, "html.parser")
                        return soup.get_text(separator="\n", strip=True)

                extracted_text = await asyncio.to_thread(_read_html)
            elif mime_type == "application/pdf":

                def _parse_pdf():
                    with pdfplumber.open(file_path) as pdf:
                        pages_text = []
                        for page in pdf.pages:
                            text = page.extract_text()
                            if text:
                                pages_text.append(text)
                        return "\n\n".join(pages_text)

                extracted_text = await asyncio.to_thread(_parse_pdf)
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                import docx

                def _parse_docx():
                    doc = docx.Document(file_path)
                    return "\n".join([p.text for p in doc.paragraphs])

                extracted_text = await asyncio.to_thread(_parse_docx)
            elif mime_type in SUPPORTED_IMAGE_TYPES:
                logger.info(f"Running OCR on {file_path}")
                try:

                    def _run_ocr():
                        with Image.open(file_path) as img:
                            return pytesseract.image_to_string(img)

                    extracted_text = await asyncio.to_thread(_run_ocr)
                except Exception as e:
                    logger.warning(f"Pytesseract failed for {file_path}: {e}")
                    return {
                        "text": "",
                        "extracted": False,
                        "error": f"OCR failed: {e}",
                        "source": TEXT_EXTRACTOR_NAME,
                    }
            elif mime_type == "text/rtf":
                from striprtf.striprtf import rtf_to_text

                def _read_rtf():
                    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                        return rtf_to_text(f.read())

                extracted_text = await asyncio.to_thread(_read_rtf)
            elif mime_type == "application/mbox":
                import contextlib
                from src.core.mbox_utils import RobustMbox

                def _parse_mbox():
                    with contextlib.closing(RobustMbox(file_path)) as mbox:
                        texts = []
                        for i, msg in enumerate(mbox):
                            if msg.is_multipart():
                                for part in msg.walk():
                                    if part.get_content_type() == "text/plain":
                                        payload = part.get_payload(decode=True)
                                        if payload:
                                            charset = (
                                                part.get_content_charset() or "utf-8"
                                            )
                                            try:
                                                text = payload.decode(
                                                    charset, errors="replace"
                                                )
                                            except LookupError:
                                                text = payload.decode(
                                                    "utf-8", errors="replace"
                                                )
                                            texts.append(text)
                            else:
                                payload = msg.get_payload(decode=True)
                                if payload:
                                    charset = msg.get_content_charset() or "utf-8"
                                    try:
                                        text = payload.decode(charset, errors="replace")
                                    except LookupError:
                                        text = payload.decode("utf-8", errors="replace")
                                    texts.append(text)
                        return "\n\n".join(texts)

                extracted_text = await asyncio.to_thread(_parse_mbox)
            elif mime_type == "application/vnd.ms-outlook":
                import extract_msg

                def _parse_msg():
                    with extract_msg.openMsg(file_path) as msg:
                        return msg.body if msg.body else ""

                extracted_text = await asyncio.to_thread(_parse_msg)
            elif mime_type == "chemical/x-cdx":
                import re

                def _read_cdx():
                    with open(file_path, "rb") as f:
                        content = f.read()
                    # Extract printable strings as a fallback for binary CDX files (ASCII range)
                    strings = re.findall(b"[\x20-\x7e]{4,}", content)
                    return "\n".join(
                        [s.decode("ascii", errors="ignore") for s in strings]
                    )

                extracted_text = await asyncio.to_thread(_read_cdx)
            elif mime_type == "application/msword":
                antiword_path = shutil.which("antiword")
                if antiword_path:
                    try:
                        # Use asyncio to run the subprocess without blocking the event loop
                        proc = await asyncio.create_subprocess_exec(
                            antiword_path,
                            "-t",
                            file_path,
                            stdout=asyncio.subprocess.PIPE,
                            stderr=asyncio.subprocess.PIPE,
                        )
                        try:
                            # Add a 30-second timeout to handle corrupt or hung documents
                            stdout, stderr = await asyncio.wait_for(
                                proc.communicate(), timeout=30.0
                            )
                            if proc.returncode == 0:
                                extracted_text = stdout.decode(
                                    "utf-8", errors="replace"
                                )
                            else:
                                logger.error(
                                    f"Antiword failed for {file_path} (code {proc.returncode}): {stderr.decode('utf-8', errors='replace')}"
                                )
                                extracted_text = ""
                        except asyncio.TimeoutError:
                            logger.error(f"Antiword TIMED OUT for {file_path}")
                            proc.kill()
                            await proc.wait()
                            extracted_text = ""
                    except Exception as e:
                        logger.error(f"Failed to run antiword for {file_path}: {e}")
                        extracted_text = ""
                else:
                    logger.warning(
                        f"Skipping {file_path}: 'antiword' is not installed. "
                        "Please install 'antiword' and ensure it is on your PATH for legacy .doc support."
                    )
                    extracted_text = ""
            elif mime_type == "application/vnd.ms-excel":
                try:

                    def _parse_excel():
                        workbook = xlrd.open_workbook(file_path)
                        all_text = []
                        for sheet in workbook.sheets():
                            all_text.append(f"Sheet: {sheet.name}")
                            for row_idx in range(sheet.nrows):
                                row_values = [
                                    str(val)
                                    for val in sheet.row_values(row_idx)
                                    if val is not None
                                ]
                                all_text.append(" ".join(row_values))
                        return "\n".join(all_text)

                    extracted_text = await asyncio.to_thread(_parse_excel)
                except Exception as e:
                    logger.error(f"xlrd failed for {file_path}: {e}")
                    extracted_text = ""
            elif mime_type in (
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ):
                try:

                    def _parse_powerpoint():
                        if (
                            mime_type
                            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        ):
                            from pptx import Presentation

                            prs = Presentation(file_path)
                            full_text = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        full_text.append(shape.text)
                            return "\n".join(full_text)
                        else:
                            parser = createParser(file_path)
                            if parser:
                                with parser:
                                    metadata = extractMetadata(parser)
                                    if metadata:
                                        plaintext = metadata.exportPlaintext()
                                        if isinstance(plaintext, str):
                                            return plaintext
                                        else:
                                            try:
                                                return "\n".join(plaintext)
                                            except TypeError:
                                                return str(plaintext)
                                    else:
                                        logger.warning(
                                            f"No metadata found for {file_path}"
                                        )
                                        return ""
                            else:
                                logger.warning(f"Hachoir could not parse {file_path}")
                                return ""

                    extracted_text = await asyncio.to_thread(_parse_powerpoint)
                except Exception as e:
                    logger.error(f"PowerPoint extraction failed for {file_path}: {e}")
                    extracted_text = ""
            elif mime_type == "message/rfc822":
                import email
                from email import policy

                try:

                    def _read_rfc822():
                        with open(file_path, "rb") as f:
                            return email.message_from_binary_file(
                                f, policy=policy.default
                            )

                    msg = await asyncio.to_thread(_read_rfc822)
                    subject = msg.get("subject", "")
                    from_addr = msg.get("from", "")
                    to_addr = msg.get("to", "")
                    date = msg.get("date", "")

                    body = ""
                    if msg.is_multipart():
                        for part in msg.walk():
                            if part.get_content_type() == "text/plain":
                                body = part.get_content()
                                break
                    else:
                        body = msg.get_content()

                    extracted_text = f"Subject: {subject}\nFrom: {from_addr}\nTo: {to_addr}\nDate: {date}\n\n{body}"
                except Exception as e:
                    logger.error(f"Email parsing failed for {file_path}: {e}")
                    extracted_text = ""
            elif mime_type == "application/vnd.wordperfect":
                import re

                try:
                    # Legacy WordPerfect files are complex; without wpd2text,
                    # we perform a robust string extraction as a fallback.
                    def _read_wp():
                        with open(file_path, "rb") as f:
                            content = f.read()
                        # Extract printable sequences of 4+ characters
                        strings = re.findall(b"[\x20-\x7e]{4,}", content)
                        extracted_text = "\n".join(
                            [s.decode("ascii", errors="ignore") for s in strings]
                        )
                        return extracted_text, len(strings)

                    extracted_text, num_strings = await asyncio.to_thread(_read_wp)
                    logger.info(
                        f"Extracted {num_strings} strings from WordPerfect file"
                    )
                except Exception as e:
                    logger.error(
                        f"WordPerfect raw extraction failed for {file_path}: {e}"
                    )
                    extracted_text = ""
            else:
                # We log the unsupported type. The extraction will be empty,
                # causing the ValueError below to mark the task as FAILED.
                logger.debug(
                    f"Skipping text extraction for unsupported mime type: {mime_type}"
                )

            extracted_content = extracted_text.strip()
            is_extracted = bool(extracted_content)

            if not is_extracted:
                if mime_type in SUPPORTED_IMAGE_TYPES:
                    logger.info(
                        f"No text extracted from image {file_path}, but skipping ValueError to prevent retries."
                    )
                else:
                    raise ValueError(
                        f"No text extracted (MIME: {mime_type or 'unknown'})"
                    )

            return {
                "text": extracted_content,
                "extracted": is_extracted,
                "source": TEXT_EXTRACTOR_NAME,
            }

        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            raise
